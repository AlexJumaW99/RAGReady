# nodes.py
import os
import sys
import json
import shutil
import datetime

import psycopg2
from psycopg2.extras import Json
from langchain_text_splitters import RecursiveCharacterTextSplitter, Language

from states import RAGIngestionState, FileMetadata
from config import (
    get_llm,
    get_genai_client,
    get_embedding_model,
    # Processing config
    OCR_METHOD,
    TRANSCRIPTION_METHOD,
    MAX_BINARY_MB,
    CHUNK_SIZE_TEXT,
    CHUNK_OVERLAP_TEXT,
    CHUNK_SIZE_CODE,
    CHUNK_OVERLAP_CODE,
    # File type sets
    TEXT_EXTENSIONS,
    CODE_EXTENSIONS,
    CONFIG_EXTENSIONS,
    PDF_EXTENSIONS,
    OFFICE_EXTENSIONS,
    STRUCTURED_EXTENSIONS,
    IMAGE_EXTENSIONS,
    AUDIO_EXTENSIONS,
    CODE_LANGUAGE_MAP,
    MIME_TYPE_MAP,
)
from utils import build_project_tree, get_file_stats


# ===========================================================================
# Private Helpers
# ===========================================================================

def _debug_error_with_llm(error_message: str, context: dict, step_name: str) -> tuple[str, str]:
    """Use Gemini to analyze an error and produce a debug summary + log file."""
    debug_prompt = f"""
An error occurred during the '{step_name}' step of a RAG ingestion pipeline.

--- CONTEXT ---
Target Path: {context.get('target_path', 'N/A')}
PostgreSQL Host: {context.get('pg_host', 'N/A')}
PostgreSQL Database: {context.get('pg_database', 'N/A')}
Last Command: {context.get('last_command', 'N/A')}
Steps Completed: {context.get('steps_completed', [])}

--- STDERR ---
{context.get('last_stderr', 'N/A')}

--- ERROR ---
{error_message}

Provide:
1. A clear explanation of the root cause.
2. Step-by-step fixes the user can apply.
"""
    try:
        response = get_llm().invoke(debug_prompt)
        summary = response.content
    except Exception as llm_err:
        summary = f"LLM debug failed: {llm_err}\n\nOriginal error:\n{error_message}"

    log_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Logs")
    os.makedirs(log_dir, exist_ok=True)
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    log_path = os.path.join(log_dir, f"error_{step_name}_{timestamp}.txt")

    with open(log_path, "w", encoding="utf-8") as f:
        f.write(f"ERROR LOG — Step: {step_name}\nTimestamp: {timestamp}\n")
        f.write("=" * 80 + "\n\n--- LLM DEBUG SUMMARY ---\n\n")
        f.write((summary or "N/A") + "\n\n--- RAW ERROR ---\n\n")
        f.write((error_message or "N/A") + "\n")

    return summary, log_path


def _build_context(state: RAGIngestionState) -> dict:
    return {
        "target_path": state.get("target_path"),
        "pg_host": state.get("pg_host"),
        "pg_database": state.get("pg_database"),
        "last_command": state.get("last_command"),
        "last_stderr": state.get("last_stderr"),
        "steps_completed": state.get("steps_completed"),
    }


def _error_return(state, step_name, error_msg, cmd_str=""):
    """Standard error‑state builder shared by every node."""
    print(f"  ❌ [{step_name}] Error: {error_msg}")
    sys.stdout.flush()
    context = _build_context(state)
    debug_summary, log_path = _debug_error_with_llm(error_msg, context, step_name)
    existing_errors = list(state.get("errors") or [])
    existing_errors.append(error_msg)
    return {
        **state,
        "current_step": step_name,
        "has_error": True,
        "errors": existing_errors,
        "error_log_path": log_path,
        "debug_summary": debug_summary,
        "last_command": cmd_str,
        "last_stderr": error_msg,
    }


def _success_step(state, step_name, updates: dict, cmd_str=""):
    """Standard success‑state builder shared by every node."""
    completed = list(state.get("steps_completed") or [])
    completed.append(step_name)
    return {
        **state,
        "current_step": step_name,
        "steps_completed": completed,
        "last_command": cmd_str,
        **updates,
    }


# ---------------------------------------------------------------------------
# Local Output Helpers
# ---------------------------------------------------------------------------

def _relative_path(filepath: str, target_root: str) -> str:
    """
    Compute the relative path of a file within the ingested target.
    This preserves the original directory structure (important for codebases).

    If target_root is a single file, returns just the filename.
    If target_root is a directory, returns the path relative to it.
    """
    abs_file = os.path.abspath(filepath)
    abs_root = os.path.abspath(target_root)

    if os.path.isfile(abs_root):
        return os.path.basename(abs_file)

    try:
        return os.path.relpath(abs_file, abs_root)
    except ValueError:
        return os.path.basename(abs_file)


def _type_to_output_subdir(file_type_or_ext: str) -> str:
    """Map a file category/extension to its output subdirectory."""
    ext = file_type_or_ext.lower().lstrip(".")
    mapping = {
        "py": "code", "js": "code", "ts": "code", "jsx": "code", "tsx": "code",
        "java": "code", "c": "code", "cpp": "code", "h": "code", "hpp": "code",
        "go": "code", "rs": "code", "rb": "code", "php": "code",
        "swift": "code", "kt": "code", "scala": "code", "sh": "code",
        "bash": "code", "r": "code", "sql": "code", "css": "code",
        "scss": "code", "less": "code", "vue": "code", "svelte": "code",
        # Text
        "md": "text", "txt": "text", "rst": "text", "html": "text",
        "htm": "text", "log": "text", "tex": "text",
        # Config
        "json": "config", "yaml": "config", "yml": "config", "toml": "config",
        "xml": "config", "env": "config", "ini": "config", "cfg": "config",
        "conf": "config",
        # PDF
        "pdf": "pdf",
        # Office
        "docx": "office", "pptx": "office",
        # Structured
        "csv": "structured", "tsv": "structured", "xlsx": "structured",
        "xls": "structured",
        # Media
        "png": "media/images", "jpg": "media/images", "jpeg": "media/images",
        "gif": "media/images", "bmp": "media/images", "tiff": "media/images",
        "webp": "media/images",
        "mp3": "media/audio", "wav": "media/audio", "m4a": "media/audio",
        "ogg": "media/audio", "flac": "media/audio", "aac": "media/audio",
        "wma": "media/audio",
    }
    return mapping.get(ext, "other")


def _save_output_file(output_dir: str, subdir: str, rel_path: str, suffix: str, content: str):
    """
    Save a processed output file locally.

    Args:
        output_dir: Root Processed_Output/<run>/ directory.
        subdir:     Type subfolder (e.g. "code", "media/images").
        rel_path:   Relative path preserving original structure (e.g. "src/main.py").
        suffix:     What to append (e.g. ".extracted.txt", ".metadata.json").
        content:    The text content to write.

    Returns:
        The full path of the written file.
    """
    out_path = os.path.join(output_dir, subdir, rel_path + suffix)
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(content)
    return out_path


def _save_json(output_dir: str, subdir: str, rel_path: str, suffix: str, data: dict | list):
    """Save a JSON output file locally (pretty-printed for review)."""
    content = json.dumps(data, indent=2, ensure_ascii=False, default=str)
    return _save_output_file(output_dir, subdir, rel_path, suffix, content)


def _copy_original(output_dir: str, subdir: str, rel_path: str, src_path: str):
    """Copy the original file into the output directory for reference."""
    out_path = os.path.join(output_dir, subdir, rel_path)
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    try:
        shutil.copy2(src_path, out_path)
    except Exception:
        pass
    return out_path


# ---------------------------------------------------------------------------
# File Classification Helper
# ---------------------------------------------------------------------------

def _classify_file(filepath: str) -> tuple[str, dict]:
    """Classify a single file into a category and return extra info."""
    ext = os.path.splitext(filepath)[1].lower()
    info = {"path": filepath, "extension": ext}

    if ext in CODE_EXTENSIONS:
        info["language"] = CODE_LANGUAGE_MAP.get(ext, "python")
        return "code", info
    if ext in TEXT_EXTENSIONS:
        return "text", info
    if ext in CONFIG_EXTENSIONS:
        return "config", info
    if ext in PDF_EXTENSIONS:
        return "pdf", info
    if ext in OFFICE_EXTENSIONS:
        info["subtype"] = ext.lstrip(".")
        return "office", info
    if ext in STRUCTURED_EXTENSIONS:
        info["subtype"] = ext.lstrip(".")
        return "structured", info
    if ext in IMAGE_EXTENSIONS:
        info["mime"] = MIME_TYPE_MAP.get(ext, "application/octet-stream")
        return "image", info
    if ext in AUDIO_EXTENSIONS:
        info["mime"] = MIME_TYPE_MAP.get(ext, "application/octet-stream")
        return "audio", info
    return "unknown", info


# ---------------------------------------------------------------------------
# Extraction Helpers
# ---------------------------------------------------------------------------

def _extract_pdf(filepath: str) -> tuple[str, str]:
    """Extract text from PDF with PyMuPDF. Falls back to Gemini Vision for image‑heavy pages."""
    import pymupdf as fitz  # PyMuPDF >= 1.24 supports this import directly

    doc = fitz.open(filepath)
    full_text = []
    method = "pymupdf"

    for page_num, page in enumerate(doc):
        text = page.get_text("text").strip()
        if len(text) < 50 and page.get_images():
            if get_genai_client():
                try:
                    pix = page.get_pixmap(dpi=200)
                    img_bytes = pix.tobytes("png")
                    from google.genai import types
                    response = get_genai_client().models.generate_content(
                        model="gemini-2.5-pro",
                        contents=[
                            types.Part.from_bytes(data=img_bytes, mime_type="image/png"),
                            "Extract ALL text from this scanned page. Return only the text.",
                        ],
                    )
                    text = response.text.strip()
                    method = "pymupdf+gemini_vision"
                except Exception:
                    pass
        full_text.append(text)

    doc.close()
    return "\n\n".join(full_text), method


def _extract_docx(filepath: str) -> str:
    """Extract text from a DOCX file."""
    from docx import Document

    doc = Document(filepath)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    return "\n\n".join(paragraphs)


def _extract_pptx(filepath: str) -> str:
    """Extract text from a PPTX file."""
    from pptx import Presentation

    prs = Presentation(filepath)
    texts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"[Slide {slide_num}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
        texts.append("\n".join(slide_texts))
    return "\n\n".join(texts)


def _ocr_image(filepath: str, method: str) -> tuple[str, str]:
    """OCR an image file using the configured method."""
    if method == "gemini" and get_genai_client():
        with open(filepath, "rb") as f:
            img_bytes = f.read()
        ext = os.path.splitext(filepath)[1].lower()
        mime = MIME_TYPE_MAP.get(ext, "image/png")
        from google.genai import types
        response = get_genai_client().models.generate_content(
            model="gemini-2.5-pro",
            contents=[
                types.Part.from_bytes(data=img_bytes, mime_type=mime),
                "Extract ALL text visible in this image. If no text is present, describe the image content in detail. Return only the extracted or described content.",
            ],
        )
        return response.text.strip(), "gemini_vision"
    else:
        try:
            import pytesseract
            from PIL import Image
            img = Image.open(filepath)
            text = pytesseract.image_to_string(img)
            return text.strip(), "tesseract"
        except ImportError:
            return "[OCR unavailable — install pytesseract: pip install pytesseract]", "none"


def _transcribe_audio(filepath: str, method: str) -> tuple[str, str]:
    """Transcribe an audio file using the configured method."""
    if method == "gemini" and get_genai_client():
        uploaded = get_genai_client().files.upload(file=filepath)
        response = get_genai_client().models.generate_content(
            model="gemini-2.5-flash",
            contents=[
                uploaded,
                "Transcribe this audio completely and accurately. Return only the transcription text.",
            ],
        )
        return response.text.strip(), "gemini_audio"
    else:
        try:
            import whisper
            model = whisper.load_model("base")
            result = model.transcribe(filepath)
            return result["text"].strip(), "whisper"
        except ImportError:
            return "[Transcription unavailable — install openai-whisper]", "none"


def _read_structured(filepath: str) -> dict:
    """Read a structured file (CSV/XLSX) and return info dict."""
    import pandas as pd

    ext = os.path.splitext(filepath)[1].lower()
    try:
        if ext in (".csv", ".tsv"):
            sep = "\t" if ext == ".tsv" else ","
            df = pd.read_csv(filepath, sep=sep, nrows=500)
        elif ext in (".xlsx", ".xls"):
            df = pd.read_excel(filepath, nrows=500)
        else:
            return {"content": "", "column_names": [], "row_count": 0}

        return {
            "content": df.to_csv(index=False),
            "column_names": list(df.columns),
            "row_count": len(df),
            "dtypes": {col: str(dtype) for col, dtype in df.dtypes.items()},
        }
    except Exception as e:
        return {"content": "", "column_names": [], "row_count": 0, "error": str(e)}


# ===========================================================================
# NODE 1: Read & Classify Files
# ===========================================================================

def read_and_classify_files(state: RAGIngestionState) -> RAGIngestionState:
    """Walks the target path, classifies every file by type, saves classification locally."""
    step_name = "read_and_classify_files"
    target_path = state.get("target_path", "")
    output_dir = state.get("output_dir", "")
    cmd_str = f"Scan & classify: {target_path}"

    print(f"\n  📂 [{step_name}] Scanning: {target_path} ...")
    sys.stdout.flush()

    try:
        classified: dict[str, list] = {
            "text": [], "code": [], "config": [], "pdf": [],
            "office": [], "structured": [], "image": [], "audio": [], "unknown": [],
        }

        if os.path.isfile(target_path):
            cat, info = _classify_file(target_path)
            classified[cat].append(info)
        elif os.path.isdir(target_path):
            for root, dirs, files in os.walk(target_path):
                dirs[:] = [d for d in dirs if not d.startswith(".") and d not in
                           {"__pycache__", "node_modules", ".git", "venv", ".venv"}]
                for fname in files:
                    if fname.startswith("."):
                        continue
                    fpath = os.path.join(root, fname)
                    cat, info = _classify_file(fpath)
                    classified[cat].append(info)
        else:
            raise FileNotFoundError(f"Path does not exist: {target_path}")

        total = sum(len(v) for v in classified.values())
        if total == 0:
            raise ValueError(f"No files found at {target_path}")

        # Build project tree
        tree = build_project_tree(target_path)

        # === SAVE LOCALLY ===
        if output_dir:
            # Save project tree
            tree_path = os.path.join(output_dir, "_project_tree.txt")
            with open(tree_path, "w", encoding="utf-8") as f:
                f.write(f"Project Tree for: {target_path}\n")
                f.write("=" * 60 + "\n\n")
                f.write(tree)

            # Save classification summary (human-readable)
            classification_summary = {
                "target_path": target_path,
                "total_files": total,
                "breakdown": {k: len(v) for k, v in classified.items() if v},
                "files_by_category": {
                    k: [{"path": f["path"], "extension": f.get("extension", "")}
                        for f in v]
                    for k, v in classified.items() if v
                },
            }
            cls_path = os.path.join(output_dir, "_classification.json")
            with open(cls_path, "w", encoding="utf-8") as f:
                json.dump(classification_summary, f, indent=2, ensure_ascii=False)

            print(f"    💾 Saved: _project_tree.txt, _classification.json")

        summary_parts = [f"{k}: {len(v)}" for k, v in classified.items() if v]
        print(f"  ✅ [{step_name}] Found {total} file(s) — {', '.join(summary_parts)}")
        sys.stdout.flush()

        return _success_step(state, step_name, {
            "classified_files": classified,
            "project_tree": tree,
            "last_stdout": f"Classified {total} files",
        }, cmd_str)

    except Exception as e:
        return _error_return(state, step_name, str(e), cmd_str)


# ===========================================================================
# NODE 2: Process Text Documents (code, text, config, PDF, office)
# ===========================================================================

def process_text_documents(state: RAGIngestionState) -> RAGIngestionState:
    """Extracts text from all text‑based files and saves extracted text locally."""
    step_name = "process_text_documents"
    classified = state.get("classified_files") or {}
    output_dir = state.get("output_dir", "")
    target_path = state.get("target_path", "")
    cmd_str = "Extract text from documents"

    text_categories = ["text", "code", "config", "pdf", "office"]
    file_lists = [classified.get(cat, []) for cat in text_categories]
    all_files = [f for sublist in file_lists for f in sublist]

    if not all_files:
        print(f"  ⏭️  [{step_name}] No text documents to process — skipping.")
        sys.stdout.flush()
        return _success_step(state, step_name, {"processed_documents": []}, cmd_str)

    print(f"  📄 [{step_name}] Processing {len(all_files)} text document(s)...")
    sys.stdout.flush()

    try:
        processed = []
        saved_count = 0

        for finfo in all_files:
            fpath = finfo["path"]
            ext = finfo.get("extension", "")
            extraction_method = "direct_read"
            extracted_text = ""

            try:
                if ext in PDF_EXTENSIONS:
                    extracted_text, extraction_method = _extract_pdf(fpath)
                elif ext == ".docx":
                    extracted_text = _extract_docx(fpath)
                    extraction_method = "python_docx"
                elif ext == ".pptx":
                    extracted_text = _extract_pptx(fpath)
                    extraction_method = "python_pptx"
                else:
                    with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                        extracted_text = f.read()
            except Exception as read_err:
                extracted_text = f"[extraction error: {read_err}]"
                extraction_method = "error"

            if not extracted_text.strip():
                extracted_text = "[empty file]"

            doc_record = {
                "filepath": fpath,
                "file_type": ext.lstrip("."),
                "language": finfo.get("language"),
                "extracted_text": extracted_text,
                "extraction_method": extraction_method,
            }
            processed.append(doc_record)

            # === SAVE LOCALLY: extracted text ===
            if output_dir:
                rel = _relative_path(fpath, target_path)
                subdir = _type_to_output_subdir(ext)
                _save_output_file(output_dir, subdir, rel, ".extracted.txt", extracted_text)
                saved_count += 1

        if saved_count:
            print(f"    💾 Saved {saved_count} .extracted.txt file(s)")
        print(f"  ✅ [{step_name}] Extracted text from {len(processed)} file(s).")
        sys.stdout.flush()

        return _success_step(state, step_name, {
            "processed_documents": processed,
            "last_stdout": f"Processed {len(processed)} text documents",
        }, cmd_str)

    except Exception as e:
        return _error_return(state, step_name, str(e), cmd_str)


# ===========================================================================
# NODE 3: Process Media Files (images → OCR, audio → transcription)
# ===========================================================================

def process_media_files(state: RAGIngestionState) -> RAGIngestionState:
    """OCRs images and transcribes audio files, saves transcripts locally."""
    step_name = "process_media_files"
    classified = state.get("classified_files") or {}
    output_dir = state.get("output_dir", "")
    target_path = state.get("target_path", "")
    cmd_str = "OCR images / transcribe audio"

    images = classified.get("image", [])
    audios = classified.get("audio", [])

    if not images and not audios:
        print(f"  ⏭️  [{step_name}] No media files to process — skipping.")
        sys.stdout.flush()
        return _success_step(state, step_name, {"processed_media": []}, cmd_str)

    print(f"  🖼️  [{step_name}] Processing {len(images)} image(s) + {len(audios)} audio file(s)...")
    sys.stdout.flush()

    try:
        processed = []

        for img_info in images:
            fpath = img_info["path"]
            try:
                transcript, method = _ocr_image(fpath, OCR_METHOD)
            except Exception as e:
                transcript, method = f"[OCR failed: {e}]", "error"

            processed.append({
                "filepath": fpath,
                "file_type": "image",
                "transcript": transcript,
                "extraction_method": method,
            })
            print(f"    🖼️  OCR ({method}): {os.path.basename(fpath)}")

            # === SAVE LOCALLY: OCR transcript + copy original ===
            if output_dir:
                rel = _relative_path(fpath, target_path)
                _save_output_file(output_dir, "media/images", rel, ".ocr.txt", transcript)
                _copy_original(output_dir, "media/images", rel, fpath)

        for aud_info in audios:
            fpath = aud_info["path"]
            try:
                transcript, method = _transcribe_audio(fpath, TRANSCRIPTION_METHOD)
            except Exception as e:
                transcript, method = f"[Transcription failed: {e}]", "error"

            processed.append({
                "filepath": fpath,
                "file_type": "audio",
                "transcript": transcript,
                "extraction_method": method,
            })
            print(f"    🎤 Transcribed ({method}): {os.path.basename(fpath)}")

            # === SAVE LOCALLY: audio transcript + copy original ===
            if output_dir:
                rel = _relative_path(fpath, target_path)
                _save_output_file(output_dir, "media/audio", rel, ".transcript.txt", transcript)
                _copy_original(output_dir, "media/audio", rel, fpath)

        if output_dir and processed:
            print(f"    💾 Saved {len(processed)} transcript(s) + originals")

        print(f"  ✅ [{step_name}] Processed {len(processed)} media file(s).")
        sys.stdout.flush()

        return _success_step(state, step_name, {
            "processed_media": processed,
            "last_stdout": f"Processed {len(processed)} media files",
        }, cmd_str)

    except Exception as e:
        return _error_return(state, step_name, str(e), cmd_str)


# ===========================================================================
# NODE 4: Process Structured Files (CSV, XLSX — kept as‑is)
# ===========================================================================

def process_structured_files(state: RAGIngestionState) -> RAGIngestionState:
    """Reads structured data files, saves preview locally."""
    step_name = "process_structured_files"
    classified = state.get("classified_files") or {}
    output_dir = state.get("output_dir", "")
    target_path = state.get("target_path", "")
    structured = classified.get("structured", [])
    cmd_str = "Read structured files"

    if not structured:
        print(f"  ⏭️  [{step_name}] No structured files — skipping.")
        sys.stdout.flush()
        return _success_step(state, step_name, {"processed_structured": []}, cmd_str)

    print(f"  📊 [{step_name}] Reading {len(structured)} structured file(s)...")
    sys.stdout.flush()

    try:
        processed = []
        for finfo in structured:
            fpath = finfo["path"]
            result = _read_structured(fpath)
            record = {
                "filepath": fpath,
                "file_type": finfo.get("subtype", "csv"),
                **result,
            }
            processed.append(record)
            print(f"    📊 Read: {os.path.basename(fpath)} — {result.get('row_count', 0)} rows, "
                  f"{len(result.get('column_names', []))} columns")

            # === SAVE LOCALLY: preview JSON + copy original ===
            if output_dir:
                rel = _relative_path(fpath, target_path)
                preview = {
                    "filepath": fpath,
                    "file_type": finfo.get("subtype", "csv"),
                    "column_names": result.get("column_names", []),
                    "dtypes": result.get("dtypes", {}),
                    "row_count": result.get("row_count", 0),
                    "first_10_rows": result.get("content", "").split("\n")[:11],  # header + 10
                }
                _save_json(output_dir, "structured", rel, ".preview.json", preview)
                _copy_original(output_dir, "structured", rel, fpath)

        if output_dir and processed:
            print(f"    💾 Saved {len(processed)} .preview.json + originals")

        print(f"  ✅ [{step_name}] Processed {len(processed)} structured file(s).")
        sys.stdout.flush()

        return _success_step(state, step_name, {
            "processed_structured": processed,
            "last_stdout": f"Processed {len(processed)} structured files",
        }, cmd_str)

    except Exception as e:
        return _error_return(state, step_name, str(e), cmd_str)


# ===========================================================================
# NODE 5: Generate Metadata (LLM‑enriched per file)
# ===========================================================================

# def generate_metadata(state: RAGIngestionState) -> RAGIngestionState:
#     """Uses Gemini to generate rich metadata. Saves per-file + combined metadata locally."""
#     step_name = "generate_metadata"
#     cmd_str = "Gemini metadata generation"
#     project_tree = state.get("project_tree", "N/A")
#     output_dir = state.get("output_dir", "")
#     target_path = state.get("target_path", "")

#     docs = state.get("processed_documents") or []
#     media = state.get("processed_media") or []
#     structured = state.get("processed_structured") or []

#     all_items = (
#         [(d["filepath"], d.get("extracted_text", "")[:6000], d.get("file_type", "")) for d in docs]
#         + [(m["filepath"], m.get("transcript", "")[:6000], m.get("file_type", "")) for m in media]
#         + [(s["filepath"], s.get("content", "")[:3000], s.get("file_type", "")) for s in structured]
#     )

#     if not all_items:
#         print(f"  ⏭️  [{step_name}] Nothing to generate metadata for.")
#         sys.stdout.flush()
#         return _success_step(state, step_name, {"file_metadata": {}}, cmd_str)

#     print(f"  🧠 [{step_name}] Generating metadata for {len(all_items)} file(s)...")
#     sys.stdout.flush()

#     try:
#         metadata_map: dict[str, dict] = {}

#         for filepath, content_sample, file_type in all_items:
#             stats = get_file_stats(filepath)

#             prompt = f"""You are an expert data cataloger. Generate comprehensive metadata for this file.

# FILE PATH: {filepath}
# FILE TYPE: {file_type}
# FILE SIZE: {stats.get('file_size_bytes', 0)} bytes
# CREATED: {stats.get('file_created_date', 'unknown')}
# MODIFIED: {stats.get('file_modified_date', 'unknown')}

# PROJECT TREE:
# {project_tree[:2000]}

# CONTENT SAMPLE:
# {content_sample}

# Return STRICTLY valid JSON (no markdown wrapping) with this structure:
# {{
#     "summary": "A detailed 2-3 sentence summary of what this file contains and its purpose.",
#     "topics": ["topic1", "topic2", "topic3"],
#     "key_entities": ["entity1", "entity2"],
#     "content_category": "code|documentation|data|media|configuration|other",
#     "quality_notes": "Brief assessment of content quality/completeness."
# }}"""

#             try:
#                 response = get_llm().invoke(prompt)
#                 raw = response.content.strip()
#                 if raw.startswith("```json"):
#                     raw = raw[7:-3].strip()
#                 elif raw.startswith("```"):
#                     raw = raw[3:-3].strip()
#                 llm_meta = json.loads(raw)
#             except Exception:
#                 llm_meta = {
#                     "summary": "Metadata extraction failed.",
#                     "topics": [],
#                     "key_entities": [],
#                     "content_category": "unknown",
#                     "quality_notes": "LLM parsing error",
#                 }

#             full_meta = {
#                 **llm_meta,
#                 **stats,
#                 "source_file_type": file_type,
#                 "source_filepath": filepath,
#                 "project_tree_snippet": project_tree[:500],
#             }
#             metadata_map[filepath] = full_meta

#             print(f"    🏷️  {os.path.basename(filepath)}: {llm_meta.get('content_category', '?')}")

#             # === SAVE LOCALLY: per-file metadata ===
#             if output_dir:
#                 rel = _relative_path(filepath, target_path)
#                 ext = os.path.splitext(filepath)[1].lower()
#                 subdir = _type_to_output_subdir(ext)
#                 _save_json(output_dir, subdir, rel, ".metadata.json", full_meta)

#         # === SAVE LOCALLY: combined metadata file ===
#         if output_dir and metadata_map:
#             combined_path = os.path.join(output_dir, "metadata", "all_metadata.json")
#             os.makedirs(os.path.dirname(combined_path), exist_ok=True)
#             with open(combined_path, "w", encoding="utf-8") as f:
#                 json.dump(metadata_map, f, indent=2, ensure_ascii=False, default=str)
#             print(f"    💾 Saved {len(metadata_map)} .metadata.json + metadata/all_metadata.json")

#         print(f"  ✅ [{step_name}] Generated metadata for {len(metadata_map)} file(s).")
#         sys.stdout.flush()

#         return _success_step(state, step_name, {
#             "file_metadata": metadata_map,
#             "last_stdout": f"Generated metadata for {len(metadata_map)} files",
#         }, cmd_str)

#     except Exception as e:
#         return _error_return(state, step_name, str(e), cmd_str)

def generate_metadata(state: RAGIngestionState) -> RAGIngestionState:
    """Uses Gemini with Structured Output to generate rich metadata."""
    step_name = "generate_metadata"
    cmd_str = "Gemini structured metadata generation"
    project_tree = state.get("project_tree", "N/A")
    output_dir = state.get("output_dir", "")
    target_path = state.get("target_path", "")

    docs = state.get("processed_documents") or []
    media = state.get("processed_media") or []
    structured = state.get("processed_structured") or []

    all_items = (
        [(d["filepath"], d.get("extracted_text", "")[:6000], d.get("file_type", "")) for d in docs]
        + [(m["filepath"], m.get("transcript", "")[:6000], m.get("file_type", "")) for m in media]
        + [(s["filepath"], s.get("content", "")[:3000], s.get("file_type", "")) for s in structured]
    )

    if not all_items:
        print(f"  ⏭️  [{step_name}] Nothing to generate metadata for.")
        sys.stdout.flush()
        return _success_step(state, step_name, {"file_metadata": {}}, cmd_str)

    print(f"  🧠 [{step_name}] Generating metadata for {len(all_items)} file(s)...")
    sys.stdout.flush()

    try:
        metadata_map: dict[str, dict] = {}
        
        # 1. Initialize the structured LLM using the TypedDict
        structured_llm = get_llm().with_structured_output(FileMetadata)

        for filepath, content_sample, file_type in all_items:
            stats = get_file_stats(filepath)

            # 2. Simplified prompt. We no longer need to explain the JSON structure 
            # because the structured_llm handles that under the hood.
            prompt = f"""You are an expert data cataloger. Generate comprehensive metadata for this file.

FILE PATH: {filepath}
FILE TYPE: {file_type}
FILE SIZE: {stats.get('file_size_bytes', 0)} bytes
CREATED: {stats.get('file_created_date', 'unknown')}
MODIFIED: {stats.get('file_modified_date', 'unknown')}

PROJECT TREE:
{project_tree[:2000]}

CONTENT SAMPLE:
{content_sample}
"""
            # 3. Direct invocation returns a dictionary directly
            try:
                llm_meta = structured_llm.invoke(prompt)
            except Exception as e:
                # This now only catches true API failures, not string parsing errors
                llm_meta = {
                    "summary": "Metadata extraction failed.",
                    "topics": [],
                    "key_entities": [],
                    "content_category": "unknown",
                    "quality_notes": f"API Error: {str(e)}",
                }

            # 4. Merge the structured LLM output with local file stats
            full_meta = {
                **llm_meta,
                **stats,
                "source_file_type": file_type,
                "source_filepath": filepath,
                "project_tree_snippet": project_tree[:500],
            }
            metadata_map[filepath] = full_meta

            print(f"    🏷️  {os.path.basename(filepath)}: {llm_meta.get('content_category', '?')}")

            # === SAVE LOCALLY: per-file metadata ===
            if output_dir:
                rel = _relative_path(filepath, target_path)
                ext = os.path.splitext(filepath)[1].lower()
                subdir = _type_to_output_subdir(ext)
                _save_json(output_dir, subdir, rel, ".metadata.json", full_meta)

        # === SAVE LOCALLY: combined metadata file ===
        if output_dir and metadata_map:
            combined_path = os.path.join(output_dir, "metadata", "all_metadata.json")
            os.makedirs(os.path.dirname(combined_path), exist_ok=True)
            with open(combined_path, "w", encoding="utf-8") as f:
                json.dump(metadata_map, f, indent=2, ensure_ascii=False, default=str)
            print(f"    💾 Saved {len(metadata_map)} .metadata.json + metadata/all_metadata.json")

        print(f"  ✅ [{step_name}] Generated metadata for {len(metadata_map)} file(s).")
        sys.stdout.flush()

        return _success_step(state, step_name, {
            "file_metadata": metadata_map,
            "last_stdout": f"Generated metadata for {len(metadata_map)} files",
        }, cmd_str)

    except Exception as e:
        return _error_return(state, step_name, str(e), cmd_str)


# ===========================================================================
# NODE 6: Setup PostgreSQL (create tables)
# ===========================================================================

_REQUIRED_TABLES = {
    "document_chunks": {
        "columns": {
            "id", "filepath", "file_type", "chunk_index", "total_chunks",
            "content", "metadata", "embedding", "created_at",
        },
        "create_sql": """
            CREATE TABLE document_chunks (
                id UUID PRIMARY KEY DEFAULT gen_random_uuid(),
                filepath TEXT NOT NULL,
                file_type TEXT NOT NULL,
                chunk_index INTEGER NOT NULL,
                total_chunks INTEGER NOT NULL,
                content TEXT NOT NULL,
                metadata JSONB,
                embedding vector(384),
                created_at TIMESTAMP WITH TIME ZONE DEFAULT NOW()
            );
        """,
    },
    "structured_files": {
        "columns": {
            "id", "filepath", "file_type", "content", "column_names",
            "row_count", "metadata", "summary_embedding", "created_at",
        },
        "create_sql": """
            CREATE TABLE structured_files (
                id UUID PRIMARY KEY DEFAULT gen_random_uuid(),
                filepath TEXT NOT NULL,
                file_type TEXT NOT NULL,
                content TEXT NOT NULL,
                column_names JSONB,
                row_count INTEGER,
                metadata JSONB,
                summary_embedding vector(384),
                created_at TIMESTAMP WITH TIME ZONE DEFAULT NOW()
            );
        """,
    },
    "media_files": {
        "columns": {
            "id", "filepath", "file_type", "binary_data", "transcript",
            "metadata", "transcript_embedding", "created_at",
        },
        "create_sql": """
            CREATE TABLE media_files (
                id UUID PRIMARY KEY DEFAULT gen_random_uuid(),
                filepath TEXT NOT NULL,
                file_type TEXT NOT NULL,
                binary_data BYTEA,
                transcript TEXT,
                metadata JSONB,
                transcript_embedding vector(384),
                created_at TIMESTAMP WITH TIME ZONE DEFAULT NOW()
            );
        """,
    },
}


def _get_existing_columns(cur, table_name: str) -> set[str]:
    """Query information_schema for the columns of an existing table."""
    cur.execute("""
        SELECT column_name
        FROM information_schema.columns
        WHERE table_schema = 'public' AND table_name = %s
    """, (table_name,))
    return {row[0] for row in cur.fetchall()}


def setup_postgres(state: RAGIngestionState) -> RAGIngestionState:
    """
    Creates (or repairs) the three RAG tables with pgvector.

    For each table:
      - If it doesn't exist → create it.
      - If it exists but has stale schema → DROP and recreate.
      - If it exists with correct schema → leave it alone.
    """
    step_name = "setup_postgres"
    cmd_str = "CREATE TABLES (document_chunks, structured_files, media_files)"

    print(f"\n  🐘 [{step_name}] Setting up PostgreSQL tables...")
    sys.stdout.flush()

    try:
        conn = psycopg2.connect(
            host=state.get("pg_host"),
            port=state.get("pg_port"),
            user=state.get("pg_username"),
            password=state.get("pg_password"),
            dbname=state.get("pg_database"),
            connect_timeout=10,
        )
        conn.autocommit = True
        cur = conn.cursor()

        cur.execute("CREATE EXTENSION IF NOT EXISTS vector;")
        print(f"    ✓ pgvector extension enabled")

        for table_name, spec in _REQUIRED_TABLES.items():
            existing_cols = _get_existing_columns(cur, table_name)

            if not existing_cols:
                cur.execute(spec["create_sql"])
                print(f"    ✓ Created table: {table_name}")
            elif not spec["columns"].issubset(existing_cols):
                missing = spec["columns"] - existing_cols
                print(f"    ⚠️  Table '{table_name}' has stale schema (missing: {missing})")
                print(f"       → Dropping and recreating '{table_name}'...")
                cur.execute(f"DROP TABLE IF EXISTS {table_name} CASCADE;")
                cur.execute(spec["create_sql"])
                print(f"    ✓ Recreated table: {table_name}")
            else:
                print(f"    ✓ Table exists with correct schema: {table_name}")

        cur.close()
        conn.close()

        print(f"  ✅ [{step_name}] Tables ready: document_chunks, structured_files, media_files")
        sys.stdout.flush()

        return _success_step(state, step_name, {
            "last_stdout": "PostgreSQL tables created/verified",
        }, cmd_str)

    except Exception as e:
        return _error_return(state, step_name, str(e), cmd_str)


# ===========================================================================
# NODE 7: Vectorize & Store
# ===========================================================================

def _get_text_splitter(language: str | None):
    """Returns the appropriate LangChain text splitter."""
    if language:
        try:
            lang_enum = Language(language)
            return RecursiveCharacterTextSplitter.from_language(
                language=lang_enum,
                chunk_size=CHUNK_SIZE_CODE,
                chunk_overlap=CHUNK_OVERLAP_CODE,
            )
        except ValueError:
            pass
    return RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE_TEXT,
        chunk_overlap=CHUNK_OVERLAP_TEXT,
        separators=["\n\n", "\n", ". ", " ", ""],
    )


def vectorize_and_store(state: RAGIngestionState) -> RAGIngestionState:
    """Chunks text, embeds, stores in PostgreSQL, and saves chunks locally."""
    step_name = "vectorize_and_store"
    cmd_str = "Vectorize & INSERT into PostgreSQL"
    metadata_map = state.get("file_metadata") or {}
    output_dir = state.get("output_dir", "")
    target_path = state.get("target_path", "")

    docs = state.get("processed_documents") or []
    media = state.get("processed_media") or []
    structured = state.get("processed_structured") or []

    print(f"\n  ⚙️  [{step_name}] Vectorizing and storing data...")
    print(f"      Documents: {len(docs)} | Media: {len(media)} | Structured: {len(structured)}")
    sys.stdout.flush()

    try:
        conn = psycopg2.connect(
            host=state.get("pg_host"),
            port=state.get("pg_port"),
            user=state.get("pg_username"),
            password=state.get("pg_password"),
            dbname=state.get("pg_database"),
        )
        cur = conn.cursor()
        total_inserted = 0
        saved_chunks = 0

        # ----- 1. Document Chunks -----
        for doc in docs:
            filepath = doc["filepath"]
            text = doc.get("extracted_text", "")
            language = doc.get("language")
            meta = metadata_map.get(filepath, {})

            if not text or text.startswith("["):
                continue

            splitter = _get_text_splitter(language)
            chunks = splitter.split_text(text)
            if not chunks:
                continue

            embeddings = get_embedding_model().encode(chunks).tolist()

            # Build local chunks data for saving
            local_chunks_data = []

            for idx, (chunk, emb) in enumerate(zip(chunks, embeddings)):
                chunk_meta = {**meta, "chunk_index": idx, "total_chunks": len(chunks)}
                formatted_emb = f"[{','.join(map(str, emb))}]"

                cur.execute("""
                    INSERT INTO document_chunks
                        (filepath, file_type, chunk_index, total_chunks, content, metadata, embedding)
                    VALUES (%s, %s, %s, %s, %s, %s, %s)
                """, (
                    filepath,
                    doc.get("file_type", "unknown"),
                    idx,
                    len(chunks),
                    chunk,
                    Json(chunk_meta),
                    formatted_emb,
                ))
                total_inserted += 1

                local_chunks_data.append({
                    "chunk_index": idx,
                    "total_chunks": len(chunks),
                    "char_count": len(chunk),
                    "content": chunk,
                })

            print(f"    📝 {os.path.basename(filepath)}: {len(chunks)} chunk(s)")

            # === SAVE LOCALLY: chunks ===
            if output_dir:
                rel = _relative_path(filepath, target_path)
                ext = os.path.splitext(filepath)[1].lower()
                subdir = _type_to_output_subdir(ext)
                _save_json(output_dir, subdir, rel, ".chunks.json", {
                    "filepath": filepath,
                    "file_type": doc.get("file_type", "unknown"),
                    "language": language,
                    "chunk_config": {
                        "chunk_size": CHUNK_SIZE_CODE if language else CHUNK_SIZE_TEXT,
                        "chunk_overlap": CHUNK_OVERLAP_CODE if language else CHUNK_OVERLAP_TEXT,
                        "splitter": f"code:{language}" if language else "recursive_text",
                    },
                    "total_chunks": len(chunks),
                    "chunks": local_chunks_data,
                })
                saved_chunks += 1

        # ----- 2. Media Files -----
        for item in media:
            filepath = item["filepath"]
            transcript = item.get("transcript", "")
            meta = metadata_map.get(filepath, {})

            binary_data = None
            try:
                file_size_mb = os.path.getsize(filepath) / (1024 * 1024)
                if file_size_mb <= MAX_BINARY_MB:
                    with open(filepath, "rb") as f:
                        binary_data = f.read()
                else:
                    meta["binary_stored"] = False
                    meta["binary_note"] = (
                        f"File too large ({file_size_mb:.1f}MB > {MAX_BINARY_MB}MB limit)."
                    )
            except Exception:
                pass

            transcript_emb = None
            if transcript and not transcript.startswith("["):
                emb = get_embedding_model().encode(transcript).tolist()
                transcript_emb = f"[{','.join(map(str, emb))}]"

            cur.execute("""
                INSERT INTO media_files
                    (filepath, file_type, binary_data, transcript, metadata, transcript_embedding)
                VALUES (%s, %s, %s, %s, %s, %s)
            """, (
                filepath,
                item.get("file_type", "unknown"),
                psycopg2.Binary(binary_data) if binary_data else None,
                transcript,
                Json(meta),
                transcript_emb,
            ))
            total_inserted += 1
            print(f"    🖼️  {os.path.basename(filepath)}: stored" +
                  (f" + transcript ({len(transcript)} chars)" if transcript else ""))

        # ----- 3. Structured Files -----
        for item in structured:
            filepath = item["filepath"]
            content = item.get("content", "")
            meta = metadata_map.get(filepath, {})
            meta["column_names"] = item.get("column_names", [])
            meta["row_count"] = item.get("row_count", 0)
            meta["dtypes"] = item.get("dtypes", {})

            summary_text = (
                f"File: {os.path.basename(filepath)}. "
                f"Columns: {', '.join(item.get('column_names', []))}. "
                f"Rows: {item.get('row_count', 0)}. "
                f"{meta.get('summary', '')}"
            )
            emb = get_embedding_model().encode(summary_text).tolist()
            summary_emb = f"[{','.join(map(str, emb))}]"

            cur.execute("""
                INSERT INTO structured_files
                    (filepath, file_type, content, column_names, row_count, metadata, summary_embedding)
                VALUES (%s, %s, %s, %s, %s, %s, %s)
            """, (
                filepath,
                item.get("file_type", "csv"),
                content,
                Json(item.get("column_names", [])),
                item.get("row_count", 0),
                Json(meta),
                summary_emb,
            ))
            total_inserted += 1
            print(f"    📊 {os.path.basename(filepath)}: {item.get('row_count', 0)} rows stored")

        conn.commit()
        cur.close()
        conn.close()

        # === SAVE LOCALLY: final overview ===
        if output_dir:
            overview = {
                "session_id": state.get("session_id"),
                "target_path": target_path,
                "ingestion_timestamp": datetime.datetime.now().isoformat(),
                "total_records_inserted": total_inserted,
                "steps_completed": list(state.get("steps_completed") or []) + [step_name],
                "summary": {
                    "documents_processed": len(docs),
                    "media_processed": len(media),
                    "structured_processed": len(structured),
                    "total_chunks_created": saved_chunks,
                },
                "postgresql": {
                    "host": state.get("pg_host"),
                    "port": state.get("pg_port"),
                    "database": state.get("pg_database"),
                    "tables": ["document_chunks", "media_files", "structured_files"],
                },
                "project_tree": state.get("project_tree", ""),
            }
            overview_path = os.path.join(output_dir, "_overview.json")
            with open(overview_path, "w", encoding="utf-8") as f:
                json.dump(overview, f, indent=2, ensure_ascii=False, default=str)

            # Write a human-readable README
            readme_path = os.path.join(output_dir, "_README.txt")
            with open(readme_path, "w", encoding="utf-8") as f:
                f.write("=" * 60 + "\n")
                f.write("  RAG Pipeline — Processed Output\n")
                f.write("=" * 60 + "\n\n")
                f.write(f"Source:     {target_path}\n")
                f.write(f"Timestamp:  {overview['ingestion_timestamp']}\n")
                f.write(f"Records:    {total_inserted} inserted into PostgreSQL\n\n")
                f.write("FOLDER STRUCTURE:\n")
                f.write("-" * 40 + "\n")
                f.write("  code/          → Extracted text, chunks, and metadata for source code files\n")
                f.write("  text/          → Extracted text, chunks, and metadata for text/docs\n")
                f.write("  config/        → Extracted text, chunks, and metadata for config files\n")
                f.write("  pdf/           → Extracted text, chunks, and metadata for PDF files\n")
                f.write("  office/        → Extracted text, chunks, and metadata for DOCX/PPTX\n")
                f.write("  media/images/  → Original images + OCR output (.ocr.txt)\n")
                f.write("  media/audio/   → Original audio + transcriptions (.transcript.txt)\n")
                f.write("  structured/    → Original CSV/XLSX + column previews (.preview.json)\n")
                f.write("  metadata/      → Combined metadata for ALL files\n\n")
                f.write("FILE SUFFIXES:\n")
                f.write("-" * 40 + "\n")
                f.write("  .extracted.txt   → Raw extracted text from the original file\n")
                f.write("  .chunks.json     → Text split into chunks (what gets vectorized)\n")
                f.write("  .metadata.json   → LLM-generated summary, topics, entities, etc.\n")
                f.write("  .ocr.txt         → OCR output from image files\n")
                f.write("  .transcript.txt  → Transcription from audio files\n")
                f.write("  .preview.json    → Column names, types, and first 10 rows\n")

            print(f"    💾 Saved _overview.json + _README.txt")

        print(f"\n  ✅ [{step_name}] Successfully stored {total_inserted} record(s) in PostgreSQL.")
        sys.stdout.flush()

        return _success_step(state, step_name, {
            "records_inserted": total_inserted,
            "last_stdout": f"Inserted {total_inserted} records",
        }, cmd_str)

    except Exception as e:
        return _error_return(state, step_name, str(e), cmd_str)
